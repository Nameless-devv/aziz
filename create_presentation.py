from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
import os


def make_nav_bar(slide, width, height, active_label="Natijalar"):
    # colors
    nav_bg = RGBColor(24, 44, 67)  # dark navy
    accent = RGBColor(255, 138, 51)  # orange

    # top bar
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(0.6)).fill.solid()
    bar = slide.shapes[-1]
    bar.fill.fore_color.rgb = nav_bg
    bar.line.fill.background()

    # nav items
    nav_items = ["Kirish", "Muammo", "Yechim", "Qanday ishlaydi", "Natijalar", "Qo'llanish", "Xulosa"]
    left = Inches(0.2)
    top = Inches(0.12)
    for item in nav_items:
        tx = slide.shapes.add_textbox(left, top, Inches(1.6), Inches(0.36))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(220, 230, 240) if item != active_label else accent
        left += Inches(1.2)


def make_results_chart(path):
    # Create a bar chart similar to attachment
    labels = ["SVD", "PCA", "Wavelet", "Autoencoder", "U-Net"]
    low = [45, 50, 48, 60, 70]
    mid = [55, 60, 62, 72, 78]
    high = [70, 68, 75, 82, 88]

    x = np.arange(len(labels))
    width = 0.25

    fig, ax = plt.subplots(figsize=(8, 4.8))
    ax.bar(x - width, low, width, label='Past shovqin', color='#dbe9f8')
    ax.bar(x, mid, width, label='O`rta shovqin', color='#6d8aa3')
    ax.bar(x + width, high, width, label='Yuoqri shovqin', color='#ff8a33')

    ax.set_ylabel('Samaradorlik (%)')
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylim(0, 100)
    ax.legend()
    ax.grid(axis='y', linestyle='--', alpha=0.4)

    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def create_presentation(out_path="Aziz_Presentation.pptx"):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Title slide
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "Aziz Robo"
    s0.placeholders[1].text = "Qishloq xo'jaligi uchun Coverage Path Planning"

    # Results slide with custom design
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_nav_bar(s, slide_width, slide_height, active_label="Natijalar")

    # Left: chart image
    chart_img = "_tmp_chart.png"
    make_results_chart(chart_img)
    left = Inches(0.6)
    top = Inches(1.0)
    pic_w = Inches(6.5)
    pic_h = Inches(4.2)
    s.shapes.add_picture(chart_img, left, top, width=pic_w, height=pic_h)

    # Right: cards
    card_left = left + pic_w + Inches(0.3)
    card_top = top
    card_w = Inches(3.0)
    card_h = Inches(0.6)

    bullet_texts = [
        ("SVD va PCA", "Eng eski matematik usullar. Ozgina yaxshi, lekin kam."),
        ("Wavelet", "Yaxshi usul. Klinikada ishlatiladi. Lekin AI dan past."),
        ("Autoencoder", "AI usuli. Wavelet dan sezilarli darajada yaxshiroq."),
        ("U-Net — G'olib!", "Barcha shovqin turlarida eng yaxshi natija. 88% samaradorlik!")
    ]

    for title, desc in bullet_texts:
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_w, card_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(242, 246, 249)
        shp.line.color.rgb = RGBColor(225, 230, 235)
        tf = shp.text_frame
        tf.margin_left = Pt(8)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(24, 44, 67)
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(80, 90, 100)
        card_top += card_h + Inches(0.18)

    # Summary dark box
    summary_top = card_top + Inches(0.1)
    summary = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, summary_top, card_w, Inches(0.8))
    summary.fill.solid()
    summary.fill.fore_color.rgb = RGBColor(28, 38, 54)
    summary.text_frame.text = "Xulosa: AI usullari (Autoencoder, U-Net) an'anaviy usullarga qaraganda yuqori samaradorlik ko'rsatadi."
    for paragraph in summary.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)

    # Cleanup
    if os.path.exists(chart_img):
        os.remove(chart_img)

    prs.save(out_path)
    print(f"Presentation saved to {out_path}")


if __name__ == '__main__':
    create_presentation()

